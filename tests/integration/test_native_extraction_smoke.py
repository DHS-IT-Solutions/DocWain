"""End-to-end smoke: dispatch_native on each format yields a non-empty ExtractionResult."""
import io

import fitz
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation

from src.extraction.adapters.dispatcher import dispatch_native


def _pdf_bytes() -> bytes:
    d = fitz.open()
    p = d.new_page()
    p.insert_text((72, 72), "Smoke PDF content for the end-to-end native extraction test.")
    buf = io.BytesIO()
    d.save(buf)
    d.close()
    return buf.getvalue()


def _docx_bytes() -> bytes:
    d = DocxDocument()
    d.add_paragraph("smoke docx content")
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _xlsx_bytes() -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "smoke"
    ws["B1"] = "xlsx"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "smoke pptx"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _csv_bytes() -> bytes:
    return b"a,b,c\n1,2,3\n"


def test_smoke_all_native_formats_produce_non_empty_output():
    cases = [
        ("smoke.pdf", _pdf_bytes(), "pdf_native"),
        ("smoke.docx", _docx_bytes(), "docx"),
        ("smoke.xlsx", _xlsx_bytes(), "xlsx"),
        ("smoke.pptx", _pptx_bytes(), "pptx"),
        ("smoke.csv", _csv_bytes(), "csv"),
    ]
    for filename, data, expected_format in cases:
        result = dispatch_native(data, filename=filename, doc_id=filename)
        assert result.format == expected_format, filename
        assert result.path_taken == "native"
        has_content = bool(result.pages) or bool(result.sheets) or bool(result.slides)
        assert has_content, f"no content extracted from {filename}"
