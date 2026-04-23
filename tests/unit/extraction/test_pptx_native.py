import io

from pptx import Presentation

from src.extraction.adapters.pptx_native import extract_pptx_native
from src.extraction.canonical_schema import ExtractionResult


def _make_pptx_two_slides_with_notes() -> bytes:
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank-ish
    s1 = prs.slides.add_slide(slide_layout)
    tx = s1.shapes.title
    if tx is not None:
        tx.text = "Slide One Title"
    s1.notes_slide.notes_text_frame.text = "notes for slide 1"

    s2 = prs.slides.add_slide(slide_layout)
    s2.notes_slide.notes_text_frame.text = "notes for slide 2"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_native_extracts_two_slides_with_notes():
    pptx = _make_pptx_two_slides_with_notes()
    result = extract_pptx_native(pptx, doc_id="d1", filename="t.pptx")
    assert isinstance(result, ExtractionResult)
    assert result.format == "pptx"
    assert len(result.slides) == 2
    assert result.slides[0].slide_num == 1
    assert result.slides[0].notes == "notes for slide 1"
    assert result.slides[1].notes == "notes for slide 2"
