"""Programmatic bench-case generator.

Creates one fixture per native format with known source bytes and known expected
JSON. Run this from project root to (re)populate tests/extraction_bench/cases/.
"""
from __future__ import annotations

import io
import json
import os
from dataclasses import asdict
from pathlib import Path

import fitz
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation

BENCH_ROOT = Path(__file__).resolve().parents[1] / "cases"


def _write_case(doc_id: str, ext: str, source_bytes: bytes, expected: dict) -> None:
    case_dir = BENCH_ROOT / doc_id
    case_dir.mkdir(parents=True, exist_ok=True)
    (case_dir / f"source{ext}").write_bytes(source_bytes)
    (case_dir / "expected.json").write_text(json.dumps(expected, indent=2), encoding="utf-8")
    (case_dir / "notes.md").write_text(f"Synthetic bench fixture for {doc_id}.\n", encoding="utf-8")


def _expected_for_pdf(text_per_page: list[str]) -> dict:
    return {
        "format": "pdf_native",
        "path_taken": "native",
        "pages": [
            {
                "page_num": i + 1,
                "blocks": [{"text": t, "block_type": "paragraph"}],
                "tables": [],
            }
            for i, t in enumerate(text_per_page)
        ],
        "sheets": [],
        "slides": [],
    }


def _expected_for_docx(paragraphs: list[str], table_rows: list[list[str]] | None) -> dict:
    blocks = [{"text": p, "block_type": "paragraph"} for p in paragraphs]
    tables = [{"rows": table_rows}] if table_rows else []
    return {
        "format": "docx",
        "path_taken": "native",
        "pages": [{"page_num": 1, "blocks": blocks, "tables": tables}],
        "sheets": [],
        "slides": [],
    }


def _expected_for_xlsx(sheet_name: str, rows: list[list]) -> dict:
    cells = {}
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row, start=1):
            cells[f"({r_idx}, {c_idx})"] = {"value": value, "formula": None, "type": "n" if isinstance(value, (int, float)) else "s"}
    return {
        "format": "xlsx",
        "path_taken": "native",
        "pages": [],
        "sheets": [{"name": sheet_name, "cells": cells, "hidden": False, "merged_cells": [], "named_ranges": []}],
        "slides": [],
    }


def _expected_for_pptx(titles: list[str], notes: list[str]) -> dict:
    slides = []
    for i, (title, note) in enumerate(zip(titles, notes), start=1):
        slides.append(
            {
                "slide_num": i,
                "elements": [{"text": title, "block_type": "paragraph"}] if title else [],
                "tables": [],
                "notes": note,
                "hidden": False,
            }
        )
    return {"format": "pptx", "path_taken": "native", "pages": [], "sheets": [], "slides": slides}


def _expected_for_csv(rows: list[list[str]]) -> dict:
    return {
        "format": "csv",
        "path_taken": "native",
        "pages": [{"page_num": 1, "blocks": [], "tables": [{"rows": rows}]}],
        "sheets": [],
        "slides": [],
    }


def generate_pdf_case():
    d = fitz.open()
    # Use longer text per page so the native-path threshold (30 chars) is satisfied.
    pages_text = ["Invoice total amount for this purchase is 1234.56 USD.", "Vendor record: Acme Corporation"]
    for t in pages_text:
        p = d.new_page()
        p.insert_text((72, 72), t)
    buf = io.BytesIO()
    d.save(buf)
    d.close()
    _write_case("bench_pdf_01", ".pdf", buf.getvalue(), _expected_for_pdf(pages_text))


def generate_docx_case():
    doc = DocxDocument()
    paras = ["First paragraph.", "Second paragraph with detail."]
    for p in paras:
        doc.add_paragraph(p)
    table_rows = [["Header A", "Header B"], ["cell 1", "cell 2"]]
    t = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
    for i, row in enumerate(table_rows):
        for j, val in enumerate(row):
            t.rows[i].cells[j].text = val
    buf = io.BytesIO()
    doc.save(buf)
    _write_case("bench_docx_01", ".docx", buf.getvalue(), _expected_for_docx(paras, table_rows))


def generate_xlsx_case():
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    ws = wb.create_sheet(title="Data")
    rows = [["Name", "Amount"], ["Alice", 100], ["Bob", 200]]
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row, start=1):
            ws.cell(row=r_idx, column=c_idx, value=value)
    buf = io.BytesIO()
    wb.save(buf)
    _write_case("bench_xlsx_01", ".xlsx", buf.getvalue(), _expected_for_xlsx("Data", rows))


def generate_pptx_case():
    prs = Presentation()
    layout = prs.slide_layouts[5]
    titles = ["Overview", "Details"]
    notes = ["speaker notes one", "speaker notes two"]
    for title, note in zip(titles, notes):
        s = prs.slides.add_slide(layout)
        if s.shapes.title is not None:
            s.shapes.title.text = title
        s.notes_slide.notes_text_frame.text = note
    buf = io.BytesIO()
    prs.save(buf)
    _write_case("bench_pptx_01", ".pptx", buf.getvalue(), _expected_for_pptx(titles, notes))


def generate_csv_case():
    rows = [["h1", "h2"], ["1", "2"], ["3", "4"]]
    data = ("\n".join(",".join(r) for r in rows) + "\n").encode("utf-8")
    _write_case("bench_csv_01", ".csv", data, _expected_for_csv(rows))


def generate_scanned_pdf_case():
    """A PDF whose pages have NO text layer — only shapes — simulating a scan."""
    d = fitz.open()
    for _ in range(1):
        p = d.new_page()
        p.draw_rect(fitz.Rect(72, 72, 400, 400), color=(0.1, 0.1, 0.1))
    buf = io.BytesIO()
    d.save(buf)
    d.close()
    expected = {
        "format": "pdf_scanned",
        "path_taken": "vision",
        "pages": [{"page_num": 1, "blocks": [{"text": "scanned content via fallback", "block_type": "paragraph"}], "tables": []}],
        "sheets": [],
        "slides": [],
    }
    _write_case("bench_scan_01", ".pdf", buf.getvalue(), expected)


def generate_image_case():
    """A PNG with a solid background — vision fallback is expected to emit a known text line."""
    from PIL import Image
    img = Image.new("RGB", (400, 200), (255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    expected = {
        "format": "image",
        "path_taken": "vision",
        "pages": [{"page_num": 1, "blocks": [{"text": "image content via fallback", "block_type": "paragraph"}], "tables": []}],
        "sheets": [],
        "slides": [],
    }
    _write_case("bench_image_01", ".png", buf.getvalue(), expected)


def main() -> None:
    BENCH_ROOT.mkdir(parents=True, exist_ok=True)
    generate_pdf_case()
    generate_docx_case()
    generate_xlsx_case()
    generate_pptx_case()
    generate_csv_case()
    generate_scanned_pdf_case()
    generate_image_case()
    print(f"generated fixtures under {BENCH_ROOT}")


if __name__ == "__main__":
    main()
