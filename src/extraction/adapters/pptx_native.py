"""Native PPTX adapter using python-pptx.

Per slide: text frames, tables, shapes with text, notes, slide masters. Hidden
slides flagged, not dropped.

Spec: docs/superpowers/specs/2026-04-23-extraction-accuracy-design.md §4.1
"""
from __future__ import annotations

import io

from pptx import Presentation

from src.extraction.adapters.errors import NativeAdapterError
from src.extraction.canonical_schema import (
    Block,
    ExtractionResult,
    Slide,
    Table,
)


def extract_pptx_native(file_bytes: bytes, *, doc_id: str, filename: str) -> ExtractionResult:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        raise NativeAdapterError(f"failed to open PPTX {filename!r}: {exc}") from exc

    slides: list[Slide] = []
    for idx, slide in enumerate(prs.slides, start=1):
        elements: list[Block] = []
        tables: list[Table] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if text:
                    elements.append(Block(text=text, bbox=None, block_type="paragraph"))
            if shape.has_table:
                t = shape.table
                rows = []
                for row in t.rows:
                    rows.append([cell.text for cell in row.cells])
                tables.append(Table(rows=rows, bbox=None, header_row_index=0 if rows else None))

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text or ""

        try:
            show = slide.element.get("show", "1")
            hidden = (show == "0")
        except Exception:
            hidden = False

        slides.append(
            Slide(
                slide_num=idx,
                elements=elements,
                tables=tables,
                images=[],
                notes=notes,
                hidden=hidden,
            )
        )

    return ExtractionResult(
        doc_id=doc_id,
        format="pptx",
        path_taken="native",
        slides=slides,
    )
